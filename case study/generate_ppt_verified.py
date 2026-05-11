import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Design Configuration ---
COLOR_HEADER = RGBColor(0, 32, 96)    # Deep Navy
COLOR_ACCENT = RGBColor(0, 112, 192)  # Professional Blue
COLOR_TEXT = RGBColor(40, 40, 40)     # Clean Slate
COLOR_WHITE = RGBColor(255, 255, 255)

FONT_MAIN = 'Arial'

def add_polished_slide(prs, title_text, sub_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header Bar
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95))
    fill.fill.solid()
    fill.fill.fore_color.rgb = COLOR_HEADER
    fill.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), prs.slide_width - Inches(1), Inches(0.45))
    tf = title.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_MAIN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Crisp Sub-heading
    if sub_text:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.52), prs.slide_width - Inches(1), Inches(0.35))
        stf = sub.text_frame
        stf.text = sub_text
        stf.paragraphs[0].font.size = Pt(14)
        stf.paragraphs[0].font.color.rgb = COLOR_ACCENT
        stf.paragraphs[0].font.bold = True
        stf.paragraphs[0].font.name = FONT_MAIN
        stf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_box(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = RGBColor(230, 230, 230)
    shape.line.width = Pt(0.5)
    return shape

def populate_bullets(shape, point_list):
    tf = shape.text_frame
    tf.word_wrap = True
    for line in point_list:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = f"● {line}"
        p.font.size = Pt(17)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_MAIN
        p.space_after = Pt(10)
        p.level = 0

def add_img(slide, image_name, left, top, width, caption):
    img_path = os.path.join('ppt', image_name)
    if os.path.exists(img_path):
        pic = slide.shapes.add_picture(img_path, left, top, width=width)
        pic.line.color.rgb = COLOR_ACCENT
        pic.line.width = Pt(1.5)
        
        cap = slide.shapes.add_textbox(left, top + pic.height + Inches(0.1), width, Inches(0.4))
        tf = cap.text_frame
        tf.word_wrap = True
        tf.text = caption
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 80, 80)
        return pic
    return None

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. TITLE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_HEADER
    
    t = slide.shapes.add_textbox(Inches(1), Inches(2.2), prs.slide_width-Inches(2), Inches(3))
    tf = t.text_frame; tf.text = "Vanguard-WSN: Verified 10x Performance"
    tf.paragraphs[0].font.size = Pt(56); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = COLOR_WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sub = tf.add_paragraph(); sub.text = "Empirical Validation of Energy-Balanced Path Trees"; sub.font.size = Pt(28); sub.alignment = PP_ALIGN.CENTER
    add_img(slide, "figure10_snapshot.png", Inches(3.6), Inches(4.5), Inches(6), "Vanguard Connectivity State (Figure 10)")

    # 2. THE PROBLEM (Figure 2)
    slide = add_polished_slide(prs, "The Energy Hole: Silent Network Burnout", "Visualizing Relayer Stress in WSN Deployment")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.3))
    populate_bullets(box, [
        "Sensors near the Base Station handle 90% higher traffic than edge nodes.",
        "Traditional rotation protocols fail to prevent 'relay hotspots'.",
        "The result is a blind-spot (Energy Hole) formed near the sink.",
        "When the center dies, the entire network becomes useless."
    ])
    add_img(slide, "figure2_deployment.png", Inches(7.1), Inches(1.6), Inches(5.7), "High-Stress Relay Zones near Base Station (Figure 2)")

    # 3. LEGACY FAILURE (Figure 5)
    slide = add_polished_slide(prs, "Existing Limits: Why Probability Fails", "The Stochastic Gambler's Fallacy in LEACH")
    box = add_content_box(slide, Inches(6.6), Inches(1.4), Inches(6.2), Inches(5.3))
    populate_bullets(box, [
        "LEACH relies on random probability to select network leaders.",
        "Picking a low-energy node as a leader causes immediate collapse.",
        "Connectivity plateaus at only 12% of physics' limit (Table 2).",
        "A deterministic brain is required for consistent longevity."
    ])
    add_img(slide, "figure5_lifetime.png", Inches(0.5), Inches(1.8), Inches(5.7), "Performance Cliff: FND Comparison Plateau (Figure 5)")

    # 4. ARCHITECTURE (Figure 1)
    slide = add_polished_slide(prs, "Centralized Oversight: The Vanguard Controller", "Absolute Visibility and Operational Integrity")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.3))
    populate_bullets(box, [
        "Software-Defined controller implemented at the Base Station.",
        "Global visibility: One 'Brain' monitoring 100% of node states.",
        "Real-time residual energy tracking to eliminate uncertainty.",
        "Deterministic Path Trees ensure zero-bottleneck data flow."
    ])
    add_img(slide, "figure1_architecture.png", Inches(7.1), Inches(1.6), Inches(5.7), "Proposed Vanguard-WSN Architecture (Figure 1)")

    # 5. ROUTING SOLUTION (Figure 3)
    slide = add_polished_slide(prs, "Optimal Data Paths: The Balanced Path Tree", "Synergy of Utility Scores and Routing Efficiency")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.5))
    populate_bullets(box, [
        "Utility Score: Selection based on [Energy] + [Connectivity].",
        "Adaptive EBPT: Routing data around fatigued sensors.",
        "Load Balancing: Mathematically shares the 'Post Office' burden.",
        "Confirmed zero-loop routing in high-density fields."
    ])
    add_img(slide, "figure3_routing_tree.png", Inches(7.1), Inches(1.8), Inches(5.7), "Constructed Energy-Balanced Path Tree (Figure 3)")

    # 6. TESTING SETUP (Table 1)
    slide = add_polished_slide(prs, "Verification: Scientific Lab Parameters", "Standardized First-Order Radio Physics Model")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(4.5))
    populate_bullets(box, [
        "Energy Model: IEEE Industry Standard First-Order Radio.",
        "Physical Scope: 100m x 100m field with 100 sentinel nodes.",
        "Software: Custom Python numerical simulation engine.",
        "Baselines: Direct comparison with LEACH results."
    ])
    add_img(slide, "table1_parameters.png", Inches(7.1), Inches(1.6), Inches(5.7), "Table 1: Physical and Simulation Parameters")

    # 7. LIFETIME PROOF (Figure 6)
    slide = add_polished_slide(prs, "Result 1: The 10.21x Stability Breakthrough", "Validated Lifecycle Extension and FND Comparison")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.8))
    populate_bullets(box, [
        "LEACH Base Death: ~97 Rounds.",
        "Vanguard Stable Life: ~993 Rounds.",
        "Achieved 1,021% extension in First Node Death (FND).",
        "Connectivity maintained for a decade of equivalent logic units."
    ])
    add_img(slide, "figure6_death_rounds.png", Inches(7.5), Inches(2), Inches(5.5), "Stability Comparison: Round 97 vs Round 993 (Figure 6)")

    # 8. DATA PROOF (Figure 8)
    slide = add_polished_slide(prs, "Result 2: Harvesting a Decade of Data", "Sustained Throughput and Information Integrity")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.8))
    populate_bullets(box, [
        "Longevity converts directly into productivity.",
        "Throughput harvest: 10,000+ packets vs LEACH ~1,000.",
        "Vanguard delivers 955% more environmental information.",
        "Consistent delivery without late-mission packet drops."
    ])
    add_img(slide, "figure8_throughput.png", Inches(7.5), Inches(2), Inches(5.5), "Accumulated Packet Delivery Proof (Figure 8)")

    # 9. FAIRNESS PROOF (Figure 7)
    slide = add_polished_slide(prs, "Result 3: Performance of Fair Selection", "Monitoring Load Distribution in Real-Time")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.8))
    populate_bullets(box, [
        "Fairness Index: Measuring how well we share the work.",
        "Utility selection ensures energy variance stays near zero.",
        "Prevents the 'Winner Takes All' energy burnout.",
        "Validated across 100+ rounds of automated monitoring."
    ])
    add_img(slide, "figure7_fairness_monitoring.png", Inches(7.5), Inches(2), Inches(5.5), "Real-Time Fairness Index Monitoring (Figure 7)")

    # 10. THERMAL PROOF (Figure 9)
    slide = add_polished_slide(prs, "Result 4: Thermal Mapping of Energy Balance", "Visual Confirmation of Zero Energy Holes")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.8))
    populate_bullets(box, [
        "Top (Legacy): Clear red hotspots where hubs worked to death.",
        "Bottom (Vanguard): Uniform green field proves balance.",
        "Evidence: The 'Smart Manager' successfully diverted hotspots.",
        "The final physical proof of the Energy Hole solution."
    ])
    add_img(slide, "figure9_heatmap.png", Inches(7.5), Inches(1.5), Inches(5.5), "Energy Thinning: Legacy vs Vanguard (Figure 9)")

    # 11. EFFICIENCY PROOF (Table 2)
    slide = add_polished_slide(prs, "Result 5: Reaching 92% of Physics' Limit", "Proximity to the Mathematical God-Line")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.8))
    populate_bullets(box, [
        "God-Line Score: The absolute theoretical limit of sensor life.",
        "Vanguard Performance: 92.1% efficiency reached.",
        "Comparison: Industry benchmarks achieve only 12% (Table 2).",
        "Conclusion: We are near the physical edge of what is possible."
    ])
    add_img(slide, "table2_performance.png", Inches(7.5), Inches(2.5), Inches(5.5), "Comparative Performance & Efficiency matrix (Table 2)")

    # 12. ABLATION PROOF (Table 3)
    slide = add_polished_slide(prs, "Proof of Logic: The Ablation Study", "Component Contribution to 10x Success")
    box = add_content_box(slide, Inches(6.6), Inches(1.4), Inches(6.2), Inches(5.3))
    populate_bullets(box, [
        "Module 1 (Utility Index): Provides 9.9x of the gain.",
        "Module 2 (EBPT Routing): Refines stability to the final 10.21x.",
        "Synergy: Combined deterministic logic is the primary driver.",
        "Validated: Both modules essential for 90%+ efficiency."
    ])
    add_img(slide, "table3_ablation.png", Inches(0.5), Inches(2.5), Inches(5.8), "Table 3: Module-by-Module Contribution Proof")

    # 13. OPTIMIZATION PROOF (Figure 4)
    slide = add_polished_slide(prs, "Balancing Acts: The Pareto Frontier", "Optimizing for Multi-Objective Performance")
    box = add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.8))
    populate_bullets(box, [
        "Trade-offs: Balancing Lifetime (Gamma) vs Transmission Speed.",
        "Visualizing the best possible 'Sweet Spots' for the network.",
        "Vanguard operates at the absolute frontier of efficiency.",
        "AI target: Reaching 99.9% of this curve automatically."
    ])
    add_img(slide, "figure4_pareto.png", Inches(7.5), Inches(2.5), Inches(5.5), "The Performance Frontier: Radar Optimization (Fig 4)")

    # 14-17 (Constraints, Future, Close, Sources)
    slide = add_polished_slide(prs, "Transparency: Constraints & Limitations", "Honest Engineering Boundaries")
    box = add_content_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5))
    populate_bullets(box, ["Radio 'tax' for control messages.", "Idealized models lack real-world multipath noise.", "Scalability for multi-kilometer fields.", "Controller dependency."])

    slide = add_polished_slide(prs, "Conclusion: A New Sentinel Standard", "Summary: Replaced Luck with Mathematical Certainty")
    box = add_content_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5))
    populate_bullets(box, ["1,021% extension validated.", "92.1% proximity to God-Line.", "Zero energy holes formed.", "Ready for IoT deployment."])

    slide = add_polished_slide(prs, "Scientific Foundations: Recognition", "Foundational Literature and Sources")
    box = add_content_box(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.8))
    populate_bullets(box, ["Heinzelman (2002): LEACH.", "Younis (2004): HEED.", "Amrita Vishwa Vidyapeetham | Dept. of Mathematics.", "Python Numerical Lab Assets."])

    prs.save('Vanguard_WSN_Final_Verified.pptx')
    print("Final Verified Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
