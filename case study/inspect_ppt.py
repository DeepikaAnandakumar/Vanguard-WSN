import os
from pptx import Presentation

def inspect_ppt(filename):
    cwd = os.getcwd()
    file_path = os.path.join(cwd, filename)
    
    if not os.path.exists(file_path):
        print(f"ERROR: File not found at {file_path}")
        return

    try:
        prs = Presentation(file_path)
        print(f"SUCCESS: Opened {filename}")
        print(f"SLIDE_COUNT: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            print(f"--- SLIDE {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Use utf-8 for safe printing
                    clean_text = shape.text.replace("\n", " | ").strip()
                    # Print using a safe encode-decode to avoid charmap errors in terminal
                    print(clean_text.encode('ascii', 'ignore').decode('ascii'))
                    
    except Exception as e:
        print(f"ERROR during inspection: {str(e)}")

if __name__ == "__main__":
    inspect_ppt("Vanguard_WSN_Final to show.pptx")
